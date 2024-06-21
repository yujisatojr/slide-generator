from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide size to widescreen 16:9
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_font(placeholder, font_size):
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            paragraph.alignment = PP_ALIGN.LEFT

def set_title_position_and_font(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.left = Inches(0.5)
    title.top = Inches(0.5)
    title.width = Inches(12)
    title.height = Inches(0.6)
    set_font(title, 32)

content_texts = {
    "slide_3": (
        "Artificial Intelligence (AI) has the potential to revolutionize the healthcare industry. "
        "AI involves the use of complex algorithms and software to emulate human cognition in the analysis of complicated medical data. "
        "The primary aim of health-related AI applications is to analyze relationships between prevention or treatment techniques and patient outcomes. "
        "AI programs are applied to practices such as diagnostics, treatment protocol development, drug development, personalized medicine, and patient monitoring and care. "
        "AI systems can be trained using data from various sources such as electronic health records (EHRs), lab results, imaging studies, and genetic profiles. "
        "These systems can identify patterns and trends that can improve patient care and provide healthcare professionals with insights that might be missed in traditional analysis."
    ),
    "slide_4": (
        "Applications of AI in diagnostics have shown remarkable advancements in recent years. "
        "Machine learning algorithms can analyze medical images such as X-rays, CT scans, and MRIs with greater accuracy and speed than human radiologists. "
        "AI can detect abnormalities in images, such as tumors or fractures, that might be overlooked by the human eye. "
        "In addition to imaging, AI is being used to analyze genetic data and predict the likelihood of developing certain diseases based on genetic markers. "
        "AI-driven diagnostic tools can assist doctors in making more accurate and timely diagnoses, leading to better patient outcomes and more efficient use of healthcare resources."
    ),
    "slide_5": (
        "AI in treatment and patient care encompasses a wide range of applications aimed at improving patient outcomes and optimizing healthcare delivery. "
        "One significant application is in personalized medicine, where AI algorithms analyze a patient’s genetic information, lifestyle, and environment to recommend tailored treatment plans. "
        "Predictive analytics powered by AI can forecast patient deterioration, readmissions, and other adverse events, allowing for proactive intervention. "
        "AI-driven tools such as virtual health assistants and chatbots provide patients with 24/7 support, answering their questions and helping them manage their conditions. "
        "Robotic process automation (RPA) can handle administrative tasks, reducing the burden on healthcare providers and allowing them to focus more on patient care."
    )
}

def set_subtitle_and_body(slide, subtitle_text, body_text):
    # Add subtitle with specific position and size
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.27), Inches(12), Inches(0.55))
    subtitle_frame = subtitle.text_frame
    subtitle_paragraph = subtitle_frame.add_paragraph()
    subtitle_paragraph.text = subtitle_text.strip()  # Remove leading and trailing whitespace
    subtitle_paragraph.font.size = Pt(16)
    subtitle_paragraph.alignment = PP_ALIGN.LEFT

    # Add body content with specific position and size
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.82), Inches(12), Inches(5.1))
    body_frame = body.text_frame
    body_frame.word_wrap = True  # Enable text wrapping
    body_paragraph = body_frame.add_paragraph()
    body_paragraph.text = body_text.strip()  # Remove leading and trailing whitespace
    body_paragraph.font.size = Pt(16)
    body_paragraph.alignment = PP_ALIGN.LEFT

# Slide 1: Cover Page
slide_1 = prs.slides.add_slide(prs.slide_layouts[6])

# Set background image with padding
padding = Inches(0.15)
background_image_path = './assets/pexels-pixabay-273250.jpg'
slide_1.shapes.add_picture(
    background_image_path, 
    padding, 
    padding, 
    prs.slide_width - 2 * padding, 
    prs.slide_height - 2 * padding
)

# Add title with specific position and size
title_1 = slide_1.shapes.add_textbox(Inches(0.49), Inches(2.72), prs.slide_width - Inches(1), Inches(1.5))
title_frame = title_1.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "Artificial Intelligence in Healthcare"
title_paragraph.font.size = Pt(44)
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
title_paragraph.alignment = PP_ALIGN.LEFT

# Add subtitle with specific position and size
subtitle_1 = slide_1.shapes.add_textbox(Inches(0.49), Inches(3.76), prs.slide_width - Inches(1), Inches(1))
subtitle_frame = subtitle_1.text_frame
subtitle_paragraph = subtitle_frame.add_paragraph()
subtitle_paragraph.text = "An Overview of AI Applications in Healthcare"
subtitle_paragraph.font.size = Pt(20)
subtitle_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
subtitle_paragraph.alignment = PP_ALIGN.LEFT

# Slide 2: Table of Contents
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout to avoid placeholder

set_title_position_and_font(slide_2, "Table of Contents")

# Add a textbox for the content with specified position and size
content_2 = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), prs.slide_height - Inches(2))
text_frame_2 = content_2.text_frame

# Add paragraphs for each item in the table of contents
toc_items = [
    "1. Introduction to AI in Healthcare",
    "2. Applications of AI in Diagnostics",
    "3. AI in Treatment and Patient Care",
    "4. Challenges and Ethical Considerations",
    "5. Future of AI in Healthcare"
]

for item in toc_items:
    p = text_frame_2.add_paragraph()
    p.text = item
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.LEFT

# Slide 3: Content One
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout
set_title_position_and_font(slide_3, "Introduction to AI in Healthcare")
subtitle_3_text = "An Overview of AI in the Medical Field"
body_3_text = content_texts["slide_3"]
set_subtitle_and_body(slide_3, subtitle_3_text, body_3_text)

# Slide 4: Content Two
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout
set_title_position_and_font(slide_4, "Applications of AI in Diagnostics")
subtitle_4_text = "How AI is Transforming Diagnostics"
body_4_text = content_texts["slide_4"]
set_subtitle_and_body(slide_4, subtitle_4_text, body_4_text)

# Slide 5: Content Three
slide_5 = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout
set_title_position_and_font(slide_5, "AI in Treatment and Patient Care")
subtitle_5_text = "Innovations in Treatment and Patient Management"
body_5_text = content_texts["slide_5"]
set_subtitle_and_body(slide_5, subtitle_5_text, body_5_text)

# Save the presentation
prs.save('AI_in_Healthcare_Presentation.pptx')
