from dotenv import load_dotenv, find_dotenv
import json
from openai import OpenAI
import os
from pptx import Presentation

load_dotenv(find_dotenv())

client = OpenAI(
    api_key=os.getenv('OPENAI_API_KEY'),
)
llm_model =  'gpt-3.5-turbo-0125'

def generate_contents(topic):
    message_text = [
        {
            "role": "system",
            "content": (
                '''
                You are a seasoned consultant specializing in creating visually appealing PowerPoint presentations.
                Based on the user-provided topic '{topic}', generate an outline for a 5-slide presentation, with each slide containing a header and content.
                Format the output in JSON as follows:
                {
                    "slides": [
                        {
                            "header": "",
                            "content": ""
                        },
                        ...
                    ]
                }
                '''
            )
        },
        {
            "role": "user",
            "content": topic
        }
    ]
    completion = client.chat.completions.create(
        model=llm_model,
        messages=message_text,
        temperature=0.7,
        max_tokens=800,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None,
        response_format={ "type": "json_object" }
    )
    return json.loads(completion.choices[0].message.content)

topic = input('Please provide a topic for the presentation: ')

result = generate_contents(topic)
slides = result["slides"]

pres = Presentation()

for slide in slides:
    slide_layout = pres.slide_layouts[1]
    new_slide = pres.slides.add_slide(slide_layout)

    if slide["header"]:
        title = new_slide.shapes.title
        title.text = slide["header"]

    if slide["content"]:
        shapes = new_slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = slide["content"]
        tf.fit_text(font_family='Arial', max_size=18, bold=False, italic=False)

pres.save("presentation.pptx")